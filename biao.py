import tkinter as tk
from tkinter import filedialog, messagebox
import pandas as pd
from pptx import Presentation
import copy
import os
from datetime import datetime


# ---------------- PPT 操作函数 ----------------

def duplicate_slide(prs, slide):
    """
    复制一页 PPT。

    prs   ：整个 PPT 对象
    slide ：需要复制的模板页
    """

    # 使用空白页版式
    blank_slide_layout = prs.slide_layouts[6]

    # 新增一页空白幻灯片
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # 把模板页中的所有元素复制到新页面
    for shape in slide.shapes:
        el = shape.element
        new_slide.shapes._spTree.insert_element_before(
            copy.deepcopy(el),
            'p:extLst'
        )

    return new_slide


def merge_runs_in_paragraph(para):
    """
    合并一个段落中的多个 run。

    PPT 里的文字有时会被拆成多个片段。
    比如看起来是 [姓名]，但实际可能被拆成 [姓] 和 [名]。
    合并后更容易替换。
    """

    runs = para.runs

    if len(runs) <= 1:
        return

    full_text = "".join(run.text for run in runs)

    runs[0].text = full_text

    for run in runs[1:]:
        run._r.getparent().remove(run._r)


def replace_text_in_shape(shape, replacements):
    """
    替换 PPT 普通文本框、形状中的文字。

    replacements 示例：
    {
        "[编号]": "20167890",
        "[姓名]": "赵杰",
        "[年龄]": "56"
    }
    """

    # 如果这个元素没有文本框，就跳过
    if not shape.has_text_frame:
        return

    # 遍历文本框中的每一个段落
    for para in shape.text_frame.paragraphs:

        # 合并 run，避免占位符被拆开导致无法替换
        merge_runs_in_paragraph(para)

        # 替换段落中的占位符
        for run in para.runs:
            for placeholder, value in replacements.items():
                if placeholder in run.text:
                    run.text = run.text.replace(
                        placeholder,
                        str(value)
                    )


def format_value(value):
    """
    格式化 Excel 单元格内容。

    作用：
    1. 空单元格不显示 nan
    2. 日期显示成 2026/5/10，而不是 2026-05-10 00:00:00
    3. 普通内容正常转成字符串
    """

    # 如果是空值，返回空字符串
    if pd.isna(value):
        return ""

    # 如果是 Excel 日期，转换成 年/月/日
    if isinstance(value, pd.Timestamp):
        return f"{value.year}/{value.month}/{value.day}"

    # 其他内容正常转字符串
    return str(value)


def generate_ppt(excel_path, ppt_path, output_folder):
    """
    根据 Excel 和 PPT 模板生成新的 PPT。

    逻辑：
    Excel 一行 = 一个人
    PPT 模板有 x 页 = 一个人的一整套模板
    Excel 有 y 行 = y 个人
    最终生成 x * y 页
    """

    # 读取 Excel 表格
    df = pd.read_excel(excel_path)

    # 打开 PPT 模板
    prs = Presentation(ppt_path)

    # 保存 PPT 模板中的所有页面
    # 例如模板有 5 页，这里就保存这 5 页
    template_slides = list(prs.slides)

    # 记录模板页数量
    template_count = len(template_slides)

    # 如果 PPT 模板没有页面，直接报错
    if template_count == 0:
        raise ValueError("PPT 模板中没有任何页面。")

    # 如果 Excel 没有数据，直接报错
    if df.empty:
        raise ValueError("Excel 中没有数据。")

    # 遍历 Excel 每一行
    # 每一行就是一个人的信息
    for index, row in df.iterrows():

        # 当前人员的替换字典
        replacements = {}

        # 遍历 Excel 的所有列名
        for col in df.columns:

            # PPT 中的占位符格式，例如 [姓名]
            placeholder = f"[{col}]"

            # 当前行当前列的值，例如 赵杰
            value = row[col]

            # 保存替换关系
            replacements[placeholder] = format_value(value)

        # 当前这个人复制完整一套 PPT 模板
        for template_slide in template_slides:

            # 复制模板中的某一页
            new_slide = duplicate_slide(prs, template_slide)

            # 替换这一页中的所有占位符
            for shape in new_slide.shapes:
                replace_text_in_shape(shape, replacements)

    # 删除原来的模板页
    # 因为前面的原始模板页只是复制源，不应该出现在最终结果中
    for _ in range(template_count):
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # 生成时间戳，避免覆盖旧文件
    timestamp = datetime.now().strftime('%Y%m%d%H%M%S')

    # 获取模板 PPT 文件名，不含后缀
    template_name = os.path.splitext(os.path.basename(ppt_path))[0]

    # 生成新文件名
    new_filename = f"{template_name}_批量生成_{timestamp}.pptx"

    # 拼接输出路径
    new_path = os.path.join(output_folder, new_filename)

    # 保存新的 PPT
    prs.save(new_path)

    return new_path


# ---------------- Tkinter GUI ----------------

def select_excel():
    """
    选择 Excel 文件。
    """

    path = filedialog.askopenfilename(
        filetypes=[("Excel Files", "*.xlsx")]
    )

    excel_entry.delete(0, tk.END)
    excel_entry.insert(0, path)


def select_ppt():
    """
    选择 PPT 模板文件。
    """

    path = filedialog.askopenfilename(
        filetypes=[("PPT Files", "*.pptx")]
    )

    ppt_entry.delete(0, tk.END)
    ppt_entry.insert(0, path)


def select_output():
    """
    选择输出文件夹。
    """

    path = filedialog.askdirectory()

    output_entry.delete(0, tk.END)
    output_entry.insert(0, path)


def generate():
    """
    点击“生成 PPT”按钮后执行。
    """

    excel_path = excel_entry.get()
    ppt_path = ppt_entry.get()
    output_folder = output_entry.get()

    if not excel_path or not ppt_path or not output_folder:
        messagebox.showwarning(
            "提示",
            "请填写 Excel 文件、PPT 模板和输出文件夹"
        )
        return

    try:
        new_path = generate_ppt(
            excel_path,
            ppt_path,
            output_folder
        )

        messagebox.showinfo(
            "完成",
            f"PPT生成完成！\n文件路径：{new_path}"
        )

    except Exception as e:
        messagebox.showerror(
            "错误",
            str(e)
        )


# ---------------- 创建窗口 ----------------

root = tk.Tk()
root.title("ppt_template_tool")


tk.Label(
    root,
    text="Excel 文件:"
).grid(
    row=0,
    column=0,
    padx=5,
    pady=5,
    sticky='e'
)

excel_entry = tk.Entry(root, width=50)
excel_entry.grid(row=0, column=1)

tk.Button(
    root,
    text="选择",
    command=select_excel
).grid(
    row=0,
    column=2,
    padx=5,
    pady=5
)


tk.Label(
    root,
    text="PPT 模板:"
).grid(
    row=1,
    column=0,
    padx=5,
    pady=5,
    sticky='e'
)

ppt_entry = tk.Entry(root, width=50)
ppt_entry.grid(row=1, column=1)

tk.Button(
    root,
    text="选择",
    command=select_ppt
).grid(
    row=1,
    column=2,
    padx=5,
    pady=5
)


tk.Label(
    root,
    text="输出文件夹:"
).grid(
    row=2,
    column=0,
    padx=5,
    pady=5,
    sticky='e'
)

output_entry = tk.Entry(root, width=50)
output_entry.grid(row=2, column=1)

tk.Button(
    root,
    text="选择",
    command=select_output
).grid(
    row=2,
    column=2,
    padx=5,
    pady=5
)


tk.Button(
    root,
    text="生成 PPT",
    command=generate,
    width=20
).grid(
    row=3,
    column=1,
    pady=10
)


root.mainloop()