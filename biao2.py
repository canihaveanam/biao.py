import tkinter as tk
from tkinter import filedialog, messagebox
import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
import os
import re
from datetime import datetime


# =========================
# 基本设置
# =========================

# 用这一列判断病人是否已经生成过
UNIQUE_COL = "住院号"

# True = Excel 空单元格不替换，PPT 里继续保留 [主诉]、[现病史] 等占位符
KEEP_BLANK_PLACEHOLDER = True

# True = 后续更新时，旧病人不重新生成页面，只补旧页面里仍然存在的 [占位符]
PATCH_EXISTING_PLACEHOLDERS = True

# 程序隐藏标记
TEMPLATE_MARKER = "[[AUTO_TEMPLATE_SLIDE]]"

PATIENT_MARKER_PREFIX = "[[AUTO_PATIENT_ID:"
PATIENT_MARKER_SUFFIX = "]]"

PAGE_MARKER_PREFIX = "[[AUTO_TEMPLATE_PAGE:"
PAGE_MARKER_SUFFIX = "]]"


# =========================
# PPT 基础操作
# =========================

def duplicate_slide(prs, slide):
    """
    复制一页 PPT。
    适合复制普通文本框、表格、形状、图片。
    如果模板里有复杂图表、SmartArt、视频等，python-pptx 复制能力可能有限。
    """
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    for shape in slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            copy.deepcopy(shape.element),
            'p:extLst'
        )

    return new_slide


def delete_slide(prs, index):
    """删除指定页"""
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    rel_id = slide_id.get(qn("r:id"))

    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def hide_slide(slide):
    """隐藏幻灯片。PPT 编辑界面仍能看到，但放映时会跳过。"""
    slide.element.set("show", "0")


def show_slide(slide):
    """显示幻灯片。"""
    slide.element.set("show", "1")


# =========================
# 文本读取
# =========================

def collect_shape_text(shape):
    """读取一个形状里的全部文本，用于扫描隐藏标记。"""
    texts = []

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            texts.append(collect_shape_text(sub_shape))

    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            texts.append(para.text)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for para in cell.text_frame.paragraphs:
                    texts.append(para.text)

    return "\n".join(t for t in texts if t)


def slide_text(slide):
    """读取一页幻灯片中的全部文字。"""
    return "\n".join(collect_shape_text(shape) for shape in slide.shapes)


def is_template_slide(slide):
    """判断是否为程序保留的隐藏模板页。"""
    return TEMPLATE_MARKER in slide_text(slide)


def get_template_slides(prs):
    """获取结果 PPT 中隐藏保存的模板页。"""
    return [slide for slide in prs.slides if is_template_slide(slide)]


def get_patient_id_from_slide(slide):
    """从幻灯片隐藏标记中读取住院号。"""
    text = slide_text(slide)

    pattern = (
        re.escape(PATIENT_MARKER_PREFIX)
        + r"(.*?)"
        + re.escape(PATIENT_MARKER_SUFFIX)
    )

    match = re.search(pattern, text)

    if match:
        return match.group(1).strip()

    return ""


def slide_has_patient_id(slide, patient_id):
    """判断某一页是否属于指定住院号。"""
    return get_patient_id_from_slide(slide) == patient_id


def get_existing_patient_ids(prs):
    """读取 PPT 中已经生成过的住院号。"""
    ids = set()

    for slide in prs.slides:
        pid = get_patient_id_from_slide(slide)
        if pid:
            ids.add(pid)

    return ids


# =========================
# 隐藏标记
# =========================

def add_hidden_marker(slide, text):
    """
    在页面外加一个很小的文本框，用作程序识别标记。
    位置在页面左上角外侧，不影响正常页面。
    """
    box = slide.shapes.add_textbox(
        Inches(-5),
        Inches(-5),
        Inches(4),
        Inches(0.25)
    )

    tf = box.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(1)


def move_template_slides_to_end(prs):
    """把隐藏模板页统一移动到最后，避免夹在病人页面中间。"""
    slide_id_list = prs.slides._sldIdLst
    slide_ids = list(slide_id_list)

    normal_ids = []
    template_ids = []

    for slide, slide_id in zip(prs.slides, slide_ids):
        if is_template_slide(slide):
            template_ids.append(slide_id)
        else:
            normal_ids.append(slide_id)

    for slide_id in slide_ids:
        slide_id_list.remove(slide_id)

    for slide_id in normal_ids + template_ids:
        slide_id_list.append(slide_id)


# =========================
# 文本替换
# =========================

def merge_runs(para):
    """
    合并文字片段，避免 [姓名] 被 PPT 拆成多个 run 后无法替换。
    """
    if len(para.runs) <= 1:
        return

    text = "".join(run.text for run in para.runs)
    para.runs[0].text = text

    for run in para.runs[1:]:
        run._r.getparent().remove(run._r)


def replace_text_frame(text_frame, replacements):
    """
    替换文本框中的占位符。
    返回替换次数。
    """
    replace_count = 0

    for para in text_frame.paragraphs:
        merge_runs(para)

        for run in para.runs:
            old_text = run.text
            new_text = old_text

            for key, value in replacements.items():
                if key in new_text:
                    replace_count += new_text.count(key)
                    new_text = new_text.replace(key, value)

            if new_text != old_text:
                run.text = new_text

    return replace_count


def replace_shape(shape, replacements):
    """
    替换普通文本框、表格、组合形状中的占位符。
    返回替换次数。
    """
    replace_count = 0

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub_shape in shape.shapes:
            replace_count += replace_shape(sub_shape, replacements)
        return replace_count

    if shape.has_text_frame:
        replace_count += replace_text_frame(shape.text_frame, replacements)

    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                replace_count += replace_text_frame(cell.text_frame, replacements)

    return replace_count


def replace_slide(slide, replacements):
    """
    替换一整页里的占位符。
    返回替换次数。
    """
    replace_count = 0

    for shape in slide.shapes:
        replace_count += replace_shape(shape, replacements)

    return replace_count


# =========================
# Excel 值处理
# =========================

def format_value(value):
    """格式化 Excel 单元格内容。"""
    if value is None:
        return ""

    value = str(value).strip()

    if value.lower() == "nan":
        return ""

    if value.endswith(" 00:00:00"):
        value = value.replace(" 00:00:00", "").replace("-", "/")

    if value.endswith(".0") and value[:-2].isdigit():
        value = value[:-2]

    return value


def row_to_replacements(row, columns):
    """
    把 Excel 一行变成：
    [列名] -> 值

    重要：
    如果 KEEP_BLANK_PLACEHOLDER=True，
    Excel 空值不会加入 replacements，
    所以 PPT 里的 [占位符] 会保留，方便以后补。
    """
    replacements = {}

    for col in columns:
        value = format_value(row[col])

        if KEEP_BLANK_PLACEHOLDER and value == "":
            continue

        replacements[f"[{col}]"] = value

    # 防止隐藏模板标记被复制到普通病人页
    replacements[TEMPLATE_MARKER] = ""

    return replacements


# =========================
# 补旧病人字段
# =========================

def patch_existing_patient(prs, patient_id, replacements):
    """
    对已经存在的病人页面做补丁：
    只替换该病人页面里仍然存在的 [占位符]。
    已经填好的正文不会动。
    图片不会动。
    删除过的页面不会恢复。
    """
    patched_slides = 0
    replace_count = 0

    for slide in prs.slides:
        if is_template_slide(slide):
            continue

        if slide_has_patient_id(slide, patient_id):
            count = replace_slide(slide, replacements)

            if count > 0:
                patched_slides += 1
                replace_count += count

    return patched_slides, replace_count


# =========================
# 核心生成 / 更新逻辑
# =========================

def generate_ppt(excel_path, ppt_path):
    """
    使用方式：

    第一次：
    原始模板 PPT + Excel
    ↓
    生成结果 PPT
    ↓
    程序在结果 PPT 最后隐藏保存一份模板页

    后续：
    上一次结果 PPT + 最新 Excel
    ↓
    旧病人：不重新生成，只补还存在的 [占位符]
    新病人：追加生成
    """

    df = pd.read_excel(
        excel_path,
        dtype=str,
        keep_default_na=False
    )

    if df.empty:
        raise ValueError("Excel 中没有数据。")

    df.columns = [str(col).strip() for col in df.columns]

    if UNIQUE_COL not in df.columns:
        raise ValueError(f"Excel 中必须有唯一识别列：{UNIQUE_COL}")

    prs = Presentation(ppt_path)

    if len(prs.slides) == 0:
        raise ValueError("PPT 中没有任何页面。")

    hidden_template_slides = get_template_slides(prs)
    existing_ids = get_existing_patient_ids(prs)

    # 有隐藏模板页，说明这是之前生成过的结果 PPT
    if hidden_template_slides:
        mode = "更新模式：追加新增病人 + 补旧病人占位符"
        template_slides = hidden_template_slides
        template_count = len(template_slides)
        first_generate = False

    # 没有隐藏模板页，说明这是原始模板 PPT
    else:
        mode = "全新生成模式：模板生成全部病人"
        template_slides = list(prs.slides)
        template_count = len(template_slides)
        first_generate = True

    appended_count = 0
    skipped_count = 0
    blank_id_count = 0
    patched_patient_count = 0
    patched_slide_count = 0
    patched_replace_count = 0

    for _, row in df.iterrows():
        patient_id = format_value(row[UNIQUE_COL])

        if not patient_id:
            blank_id_count += 1
            continue

        replacements = row_to_replacements(row, df.columns)

        # 如果这个住院号已经存在，说明是旧病人
        if patient_id in existing_ids:
            skipped_count += 1

            if PATCH_EXISTING_PLACEHOLDERS:
                slides_count, replace_count = patch_existing_patient(
                    prs,
                    patient_id,
                    replacements
                )

                if replace_count > 0:
                    patched_patient_count += 1
                    patched_slide_count += slides_count
                    patched_replace_count += replace_count

            continue

        # 新病人：复制模板页，追加到 PPT 后面
        for page_index, template_slide in enumerate(template_slides, start=1):
            new_slide = duplicate_slide(prs, template_slide)
            show_slide(new_slide)

            replace_slide(new_slide, replacements)

            # 每一页都写隐藏标记：住院号 + 模板页序号
            add_hidden_marker(
                new_slide,
                f"{PATIENT_MARKER_PREFIX}{patient_id}{PATIENT_MARKER_SUFFIX}"
            )

            add_hidden_marker(
                new_slide,
                f"{PAGE_MARKER_PREFIX}{page_index}{PAGE_MARKER_SUFFIX}"
            )

        existing_ids.add(patient_id)
        appended_count += 1

    # 第一次生成时：
    # 把原始模板页复制一份隐藏保存到结果 PPT 最后，
    # 然后删除最前面的原始模板页。
    if first_generate:
        original_template_slides = template_slides

        for template_slide in original_template_slides:
            hidden_copy = duplicate_slide(prs, template_slide)
            add_hidden_marker(hidden_copy, TEMPLATE_MARKER)
            hide_slide(hidden_copy)

        # 删除最前面的原始模板页，只保留病人页 + 最后的隐藏模板页
        for _ in range(template_count):
            delete_slide(prs, 0)

    # 每次保存前，把隐藏模板页挪到最后
    move_template_slides_to_end(prs)

    folder = os.path.dirname(ppt_path)
    name = os.path.splitext(os.path.basename(ppt_path))[0]

    name = re.sub(r'(_(?:批量生成|追加更新|追加更新补字段|当前版|最新版|更新版)_\d{8}_\d{6})+$', '', name)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

    new_path = os.path.join(
        folder,
        f"{name}_追加更新补字段_{timestamp}.pptx"
    )

    prs.save(new_path)

    return {
        "new_path": new_path,
        "mode": mode,
        "excel_rows": len(df),
        "template_count": template_count,
        "appended_count": appended_count,
        "skipped_count": skipped_count,
        "blank_id_count": blank_id_count,
        "generated_pages": appended_count * template_count,
        "patched_patient_count": patched_patient_count,
        "patched_slide_count": patched_slide_count,
        "patched_replace_count": patched_replace_count
    }


# =========================
# Tkinter 界面
# =========================

def select_excel():
    path = filedialog.askopenfilename(
        title="选择 Excel 文件",
        filetypes=[("Excel Files", "*.xlsx")]
    )

    if path:
        excel_entry.delete(0, tk.END)
        excel_entry.insert(0, path)


def select_ppt():
    path = filedialog.askopenfilename(
        title="选择 PPT 模板或上一次结果 PPT",
        filetypes=[("PowerPoint Files", "*.pptx")]
    )

    if path:
        ppt_entry.delete(0, tk.END)
        ppt_entry.insert(0, path)


def generate():
    excel_path = excel_entry.get().strip()
    ppt_path = ppt_entry.get().strip()

    if not excel_path or not ppt_path:
        messagebox.showwarning(
            "提示",
            "请选择 Excel 文件和 PPT 文件。"
        )
        return

    try:
        result = generate_ppt(
            excel_path,
            ppt_path
        )

        os.startfile(result["new_path"])

        messagebox.showinfo(
            "完成",
            f"PPT 处理完成！\n\n"
            f"本次模式：{result['mode']}\n"
            f"Excel 数据：{result['excel_rows']} 行\n"
            f"模板页数：{result['template_count']} 页\n\n"
            f"新增病人：{result['appended_count']} 人\n"
            f"跳过旧病人：{result['skipped_count']} 人\n"
            f"住院号为空：{result['blank_id_count']} 行\n"
            f"本次新增页数：{result['generated_pages']} 页\n\n"
            f"补字段病人：{result['patched_patient_count']} 人\n"
            f"补字段页面：{result['patched_slide_count']} 页\n"
            f"补字段替换次数：{result['patched_replace_count']} 处\n\n"
            f"文件位置：\n{result['new_path']}\n\n"
            f"注意：\n"
            f"1. 结果 PPT 最后有隐藏模板页，不要删除。\n"
            f"2. 后续更新要选择上一次结果 PPT。\n"
            
        )

    except Exception as e:
        messagebox.showerror(
            "错误",
            str(e)
        )


root = tk.Tk()
root.title("PPTBatch")
root.resizable(False, False)


tk.Label(root, text="Excel 文件：").grid(
    row=0,
    column=0,
    padx=8,
    pady=8,
    sticky="e"
)

excel_entry = tk.Entry(root, width=60)
excel_entry.grid(row=0, column=1, padx=5, pady=8)

tk.Button(root, text="选择", command=select_excel).grid(
    row=0,
    column=2,
    padx=8,
    pady=8
)


tk.Label(root, text="PPT模板/上次结果：").grid(
    row=1,
    column=0,
    padx=8,
    pady=8,
    sticky="e"
)

ppt_entry = tk.Entry(root, width=60)
ppt_entry.grid(row=1, column=1, padx=5, pady=8)

tk.Button(root, text="选择", command=select_ppt).grid(
    row=1,
    column=2,
    padx=8,
    pady=8
)


tk.Button(
    root,
    text="生成 / 追加 / 补字段",
    command=generate,
    width=26
).grid(
    row=2,
    column=1,
    pady=15
)


root.mainloop()